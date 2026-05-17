"""Regenerate ontoink.pptx for the 0.6.0 release.

Run from the repo root:

    python scripts/generate_pptx.py

Produces ``ontoink.pptx`` (overwrites existing).
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Theme ────────────────────────────────────────────────────────────────────
PRIMARY = RGBColor(0x08, 0x91, 0xb2)   # teal
ACCENT  = RGBColor(0x0e, 0x73, 0x90)
TEXT    = RGBColor(0x1f, 0x29, 0x37)
SUBTLE  = RGBColor(0x6b, 0x72, 0x80)
LIGHT   = RGBColor(0xe5, 0xe7, 0xeb)
BG      = RGBColor(0xf9, 0xfa, 0xfb)
INFER   = RGBColor(0xa8, 0x55, 0xf7)


def add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[0]   # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(54)
            run.font.bold = True
            run.font.color.rgb = PRIMARY
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle
        for para in slide.placeholders[1].text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(22)
                run.font.color.rgb = SUBTLE
    return slide


def add_content_slide(prs, title, paragraphs):
    """Add a 'Title and Content' slide.

    ``paragraphs`` is a list of (text, level) tuples. Level 0 = top-level bullet,
    higher levels indent. A leading "## " marks a sub-header.
    """
    layout = prs.slide_layouts[1]   # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = PRIMARY

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for text, level in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if text.startswith("## "):
            p.text = text[3:]
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.color.rgb = ACCENT
            p.level = 0
        else:
            p.text = text
            p.level = level
            for run in p.runs:
                run.font.size = Pt(15 if level == 0 else 13)
                run.font.color.rgb = TEXT
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title ─────────────────────────────────────────────────────
    add_title_slide(
        prs,
        "ontoink v0.6.0",
        "Interactive Ontology Visualization, SHACL Validation,\n"
        "OWL Reasoning, and Inline Editing — as Living Documentation",
    )

    # ── Slide 2: What is ontoink? ──────────────────────────────────────────
    add_content_slide(prs, "What is ontoink?", [
        ("An open-source MkDocs plugin and standalone playground that turns Turtle ontology files into interactive, publication-ready diagrams — embedded directly in Markdown documentation.", 0),
        ("## Five activities, one workflow", 0),
        ("Visualization: formal notation, six layouts, color-coded ontology sources", 0),
        ("SHACL validation: constraint overlay with cardinality badges, live re-validation", 0),
        ("OWL reasoning: four selectable backends (HermiT, Konclude native, Konclude WASM, OWL-RL)", 0),
        ("Inline Turtle editing: edit → re-parse → re-validate without reload", 0),
        ("Export: high-resolution PNG / SVG with embedded legend and prefixes", 0),
        ("## Available as", 0),
        ("MkDocs plugin (Python, pip-installable)", 0),
        ("Browser-only Playground (paste TTL → visualize, no installation)", 0),
        ("Production Docker image (docs server + REST API + reasoning service)", 0),
    ])

    # ── Slide 3: Why ontoink? ──────────────────────────────────────────────
    add_content_slide(prs, "Why ontoink?", [
        ("## The problem", 0),
        ("Ontology documentation, visualization, and validation are usually done with separate, disconnected tools.", 0),
        ("Existing visualization tools (WebVOWL, Graffoo, Chowlk) produce static diagrams that drift from the source.", 0),
        ("General-purpose diagram plugins for MkDocs (D2, Mermaid) are not ontology-aware: they cannot parse RDF, dereference IRIs, overlay SHACL, or run OWL reasoning.", 0),
        ("## ontoink's answer", 0),
        ("Treats the Turtle source as the single point of truth.", 0),
        ("All five activities run from one source file inside the documentation page.", 0),
        ("Click any node to dereference its IRI through content negotiation — labels and comments come back live from the source ontology.", 0),
        ("Browser-side OWL-DL reasoning runs entirely client-side (a first in this class of tool).", 0),
    ])

    # ── Slide 4: Architecture ──────────────────────────────────────────────
    add_content_slide(prs, "Architecture", [
        ("## Build-time (Python)", 0),
        ("rdflib parses Turtle, classifies nodes (Class / Individual / Literal / Datatype)", 0),
        ("pySHACL validates the data graph against the shape graph", 0),
        ("Optional pre-computed inferences baked into the page (HermiT or OWL-RL)", 0),
        ("Output: a single JSON blob embedded in the HTML, Base64-encoded", 0),
        ("## Runtime (JavaScript, in the browser)", 0),
        ("Cytoscape.js renders the graph; dagre / cose / circle / concentric / tree / grid layouts", 0),
        ("CodeMirror provides the Turtle editor with syntax highlighting", 0),
        ("Browser-side parser handles the playground (no server round-trip)", 0),
        ("Web Worker hosts Konclude WebAssembly for in-browser OWL-DL reasoning", 0),
        ("## Production (Docker image)", 0),
        ("ONTOINK_MODE=serve / build / api / all — one image, four roles", 0),
        ("/reason and /validate FastAPI endpoints with per-request reasoner override", 0),
    ])

    # ── Slide 5: Interactive Graph Features ────────────────────────────────
    add_content_slide(prs, "Interactive Graph Features", [
        ("Six layouts: dagre, cose (force), circle, concentric, tree, grid — switchable live", 0),
        ("Color-coded by ontology source (BFO magenta, IAO orange, NFDIcore purple, PMD teal, FOAF, Schema.org)", 0),
        ("Minimap for navigation, fuzzy search with highlight, fullscreen view", 0),
        ("Click popups with copy-IRI / copy-label and a More… button that fetches labels and axioms from the source", 0),
        ("## New in 0.6.0", 0),
        ("Edit Layout — change colors, shapes, line styles, arrow shapes per type or per namespace", 0),
        ("Legend and namespace overlays refresh live when the layout or any style changes", 0),
        ("PNG and SVG exports include all Edit-Layout customizations in the embedded legend", 0),
        ("SHACL constraint edges render in the browser playground (not only in build-time output)", 0),
        ("Labels and axioms auto-fetched for IRIs referenced in sh:path, sh:targetClass, sh:class, sh:datatype, sh:node", 0),
    ])

    # ── Slide 6: SHACL Validation & Editor ─────────────────────────────────
    add_content_slide(prs, "SHACL Validation & Editor", [
        ("SHACL constraints rendered as bold cyan dashed edges with [min..max] cardinality badges directly on the graph", 0),
        ("Build-time validation via pySHACL — violations shown with focus node, property path, severity, message", 0),
        ("Inline Turtle editor (CodeMirror) for live data editing", 0),
        ("## Edit → Validate → Visualize loop (client-side)", 0),
        ("Edit Turtle in the page, click Validate — pySHACL is not needed; a lightweight JS validator runs in the browser", 0),
        ("Click Update Graph to re-parse and re-render without page reload", 0),
        ("Validate with Inferences re-runs validation over the union of asserted and inferred triples", 0),
        ("## Standalone SHACL Editor (browser-only)", 0),
        ("Build SHACL shapes visually with templates, tooltips, live preview, class/property dropdowns", 0),
        ("Powered by the Shape Recommender (next slide)", 0),
    ])

    # ── Slide 7: Shape Recommender — Detailed ─────────────────────────────
    add_content_slide(prs, "Shape Recommender — Data Profiling", [
        ("Automatically suggests SHACL shapes from instance data or a SPARQL endpoint, based on Mihindukulasooriya et al. (2018) RDF Shape Induction.", 0),
        ("## Algorithm", 0),
        ("1. Class discovery — collect every IRI used as object of rdf:type", 0),
        ("2. Property profiling — for each class, count how often each property is used per instance", 0),
        ("3. Constraint inference — translate frequencies into SHACL:", 0),
        ("   ≥ 90% of instances have property P  →  sh:minCount 1", 1),
        ("   no instance has more than one value  →  sh:maxCount 1", 1),
        ("   all values share an XSD datatype     →  sh:datatype", 1),
        ("   all values are IRIs of class C       →  sh:class C and sh:nodeKind sh:IRI", 1),
        ("4. Confidence — percentage of instances exhibiting the pattern (0–100 %)", 0),
        ("## Extensions beyond Mihindukulasooriya (2018)", 0),
        ("sh:pattern — auto-detects emails, URLs, ALL-CAPS codes", 0),
        ("sh:minLength / sh:maxLength — string length statistics", 0),
        ("sh:minInclusive / sh:maxInclusive — numeric range constraints", 0),
        ("Uniqueness detection — flags potential identifiers", 0),
    ])

    # ── Slide 8: Shape Recommender — Worked Example ───────────────────────
    add_content_slide(prs, "Shape Recommender — Worked Example", [
        ("## Input: 100 instances of ex:Person", 0),
        ("98 / 100 have a foaf:name (xsd:string, length 4–62)", 0),
        ("100 / 100 have an ex:age (xsd:integer, range 18–95)", 0),
        ("95 / 100 have an ex:email matching ^.+@.+\\..+$", 0),
        ("60 / 100 have one or more foaf:knows pointing to another ex:Person", 0),
        ("## Recommended shape (auto-generated)", 0),
        ("ex:PersonShape a sh:NodeShape ; sh:targetClass ex:Person ;", 0),
        ("  sh:property [ sh:path foaf:name ; sh:datatype xsd:string ;", 1),
        ("                sh:minCount 1 ; sh:maxCount 1 ;", 2),
        ("                sh:minLength 4 ; sh:maxLength 62 ] ;     # 98 % confidence", 2),
        ("  sh:property [ sh:path ex:age ; sh:datatype xsd:integer ;", 1),
        ("                sh:minCount 1 ; sh:maxCount 1 ;", 2),
        ("                sh:minInclusive 18 ; sh:maxInclusive 95 ] ;   # 100 % confidence", 2),
        ("  sh:property [ sh:path ex:email ; sh:datatype xsd:string ;", 1),
        ("                sh:minCount 1 ; sh:maxCount 1 ;", 2),
        ("                sh:pattern '^.+@.+\\\\..+$' ] ;                 # 95 % confidence", 2),
        ("  sh:property [ sh:path foaf:knows ; sh:class ex:Person ;", 1),
        ("                sh:nodeKind sh:IRI ] .                          # 60 % confidence", 2),
        ("## Workflow", 0),
        ("Navigate suggestions with Prev / Next, edit, download, or accept into the visual editor.", 0),
        ("Each shape carries a confidence score so the user can filter low-evidence patterns.", 0),
    ])

    # ── Slide 9: Shape Recommender — Future TODOs ─────────────────────────
    add_content_slide(prs, "Shape Recommender — Future Work", [
        ("## In progress", 0),
        ("Anti-pattern integration — feed OntoSniff signals (lazy class, missing label) into suggestions", 0),
        ("Disjunctive paths — sh:or for properties whose values fall into clearly distinct sub-groups", 0),
        ("## Planned", 0),
        ("Property co-occurrence — recommend sh:and / sh:not based on which properties appear together", 0),
        ("Cardinality from cardinality classes — detect OWL functional / inverse-functional properties and translate to sh:maxCount 1", 0),
        ("Property paths — extend to sequence and alternative paths (sh:path with rdf:List)", 0),
        ("Active learning — let the user accept / reject individual constraints; weight future suggestions accordingly", 0),
        ("Cross-class shape merging — find common patterns across sibling classes and lift them to the superclass", 0),
        ("Sample-size warnings — emphasise low evidence when fewer than 30 instances support a constraint", 0),
        ("Integration with the LOV catalog to favour widely-used predicates and datatypes", 0),
        ("## Research questions", 0),
        ("How to balance recall (don't miss constraints) and precision (don't over-constrain) when the data is noisy?", 0),
        ("Can large-language-model context help suggest sh:message text in the user's domain language?", 0),
    ])

    # ── Slide 10: OWL Reasoning — Backends ────────────────────────────────
    add_content_slide(prs, "OWL Reasoning — Four Backends", [
        ("Choose the right reasoner for each setting. ONTOINK_REASONER selects globally; the playground dropdown selects per request.", 0),
        ("## owlready2 (HermiT)", 0),
        ("Full OWL DL via the HermiT tableau reasoner", 0),
        ("Requires Java; included in the production Docker image", 0),
        ("Best general-purpose backend for richer instance-level inferences", 0),
        ("## Konclude (native C++ binary)", 0),
        ("Upstream Konclude from University of Ulm (Steigmiller, Liebig, Glimm 2014)", 0),
        ("Optimized SROIQ tableau reasoner; runs without Java", 0),
        ("Expects OWL/XML input — best for proper OWL ontology files", 0),
        ("## Konclude WASM (rdf-reasoner-konclude)", 0),
        ("Konclude compiled to WebAssembly by Thomas Hanke", 0),
        ("Runs entirely in the browser via a Web Worker — no server needed", 0),
        ("Requires page to be cross-origin isolated (a vendored service worker handles this on static hosts)", 0),
        ("## OWL-RL (pure Python)", 0),
        ("Lightweight rules-based fallback for environments without Java or a native binary", 0),
        ("Fastest startup; covers the OWL-RL profile", 0),
    ])

    # ── Slide 11: Reasoning UX ────────────────────────────────────────────
    add_content_slide(prs, "Reasoning — Interactive Panel", [
        ("Every diagram (the playground and every fence-rendered graph) has the same interactive Reasoning panel.", 0),
        ("## Backend selector", 0),
        ("Dropdown lists every reasoner — disabled options carry a hint (\"needs cross-origin isolation\", \"server offline\")", 0),
        ("Defaults: Browser Konclude WASM when isolated, Server otherwise", 0),
        ("Each request reports which reasoner was actually used", 0),
        ("## Run state", 0),
        ("Spinner + Stop button while running; second clicks are ignored; AbortController cancels server fetches", 0),
        ("Reasoning button and dropdown disabled until the request completes or is cancelled", 0),
        ("## Result panel", 0),
        ("Stats row: count · elapsed ms · distinct subjects · distinct predicates · backend", 0),
        ("Inferred triples in a sortable table; literals carry a badge", 0),
        ("Show inferences on graph — toggles a purple-dotted overlay; color and line style adjustable in Edit Layout", 0),
        ("Download as N-Triples, copy as JSON, Re-run with a different backend", 0),
        ("Auto-expanding log section with per-step timestamps for debugging", 0),
    ])

    # ── Slide 12: Production Deployment ──────────────────────────────────
    add_content_slide(prs, "Production Deployment", [
        ("One Docker image, four roles selected by ONTOINK_MODE.", 0),
        ("## serve", 0),
        ("MkDocs dev server — docs only, no reasoning API", 0),
        ("## build", 0),
        ("One-shot static site build, useful for CI / static export", 0),
        ("## api", 0),
        ("FastAPI endpoints: /reason, /validate, /health — no docs", 0),
        ("## all  (recommended)", 0),
        ("Builds the docs once at startup, then FastAPI serves them at / and the API at /reason and /health", 0),
        ("Docs and reasoning share the same origin — the playground's Server reasoner option works without a reverse proxy", 0),
        ("## Persistence", 0),
        ("Each /reason call writes input.ttl, inferences.json, and inferences.nt to a mounted volume", 0),
        ("One time-stamped subdirectory per request; the run id is returned in the response", 0),
        ("## Reasoner override", 0),
        ("ONTOINK_REASONER env var, or a reasoner field in the JSON request body — no restart needed", 0),
    ])

    # ── Slide 13: OntoSniff ──────────────────────────────────────────────
    add_content_slide(prs, "OntoSniff — Ontology Smell Detector", [
        ("Anti-pattern detector backed by published ontology-quality research (Poveda-Villalón et al. 2014, Rector et al. 2004, Gangemi et al. 2006).", 0),
        ("## Detectors", 0),
        ("Lazy Class, Missing Label, Missing Domain / Range, Singleton Hierarchy, Property Soup, Orphan Class, Missing Inverse, No SHACL Coverage, Label Language Gap, and more", 0),
        ("Each smell has id, name, severity (info / warning / error), description, fix hints, and an example", 0),
        ("## Output", 0),
        ("Quality score 0–100 with severity breakdown", 0),
        ("Annotated report on a standalone OntoSniff page", 0),
        ("Shareable via URL parameter: ?source=https://example.org/ontology.ttl", 0),
        ("Integrated into the Stats panel for every diagram, so smells appear next to size and depth metrics", 0),
    ])

    # ── Slide 14: SPARQL & Browser Tools ─────────────────────────────────
    add_content_slide(prs, "SPARQL & Browser-Only Tools", [
        ("## Built-in SPARQL panel", 0),
        ("Query the loaded graph with SELECT; highlight matching nodes back on the canvas", 0),
        ("Templates, class and property dropdowns, autocomplete: Ctrl+Space on Windows / Linux, Alt-/ (Option-/) on macOS", 0),
        ("## SPARQL Explorer (standalone page)", 0),
        ("Connect to any SPARQL endpoint; auto-discover classes, properties, and labels", 0),
        ("Adaptive discovery for million-triple endpoints (DBpedia, Wikidata)", 0),
        ("Ontology source fallback when the endpoint has no rdfs:label — fetches BFO, IAO, NFDIcore, FOAF, etc., and extracts labels client-side", 0),
        ("## Playground", 0),
        ("Paste, upload, or fetch a Turtle file → visualize, edit, validate, reason — all in the browser", 0),
        ("Service worker enables WebAssembly Konclude on GitHub Pages without server-side header support", 0),
        ("## SHACL Editor", 0),
        ("Visual SHACL authoring with the Shape Recommender baked in", 0),
    ])

    # ── Slide 15: Tech Stack & Testing ──────────────────────────────────
    add_content_slide(prs, "Technology Stack & Testing", [
        ("## Server / build", 0),
        ("Python ≥ 3.9, MkDocs ≥ 1.4, rdflib, pySHACL, pymdown-extensions", 0),
        ("Optional: owlready2 (HermiT), FastAPI + uvicorn (API mode)", 0),
        ("Docker base: python:3.12-slim + Node 20 + Java JRE + native Konclude binary", 0),
        ("## Browser", 0),
        ("Cytoscape.js (graph), dagre (layout), cytoscape-svg (vector export)", 0),
        ("CodeMirror 5 with Turtle mode", 0),
        ("rdf-reasoner-konclude + n3 (lazy-loaded via esm.sh)", 0),
        ("coi-serviceworker vendored for cross-origin isolation on static hosts", 0),
        ("## Tests", 0),
        ("pytest: parser, label resolution, OntoSniff, SHACL extraction, reasoning dispatch, every API endpoint — 52 passing", 0),
        ("node --check on ontoink.js; manual playground checks per release", 0),
        ("GitHub Actions: CI on push / PR, Docker image build on every v* tag", 0),
    ])

    # ── Slide 16: Citations & Acknowledgments ───────────────────────────
    add_content_slide(prs, "Citations & Acknowledgments", [
        ("## Reasoners", 0),
        ("Glimm, B., Horrocks, I., Motik, B., Stoilos, G., Wang, Z. (2014). HermiT: An OWL 2 Reasoner. Journal of Automated Reasoning 53(3).", 0),
        ("Steigmiller, A., Liebig, T., Glimm, B. (2014). Konclude: System Description. Journal of Web Semantics 27–28, 78–85.", 0),
        ("Hanke, T. (2026). rdf-reasoner-konclude — OWL-DL tableau reasoning for browsers and Node.js via WebAssembly.", 0),
        ("W3C OWL 2 Web Ontology Language Profiles (OWL-RL Recommendation).", 0),
        ("## Foundations", 0),
        ("Mihindukulasooriya, N. et al. (2018). RDF Shape Induction using Knowledge Base Profiling.", 0),
        ("Poveda-Villalón, M. et al. (2014). OOPS! (OntOlogy Pitfall Scanner!).", 0),
        ("rdflib (Krech 2023) · pySHACL (Car 2020) · owlready2 (Lamy 2017)", 0),
        ("## Funding", 0),
        ("German Research Foundation (DFG) — NFDI-MatWerk (DFG 460247524) and NFDI4Culture (DFG 441958017).", 0),
    ])

    # ── Slide 17: Links & Thank You ─────────────────────────────────────
    add_content_slide(prs, "Links & Thank You", [
        ("Live demo:  https://ise-fizkarlsruhe.github.io/ontoink/", 0),
        ("Playground: https://ise-fizkarlsruhe.github.io/ontoink/playground/", 0),
        ("Source:     https://github.com/ISE-FIZKarlsruhe/ontoink", 0),
        ("PyPI:       https://pypi.org/project/ontoink/", 0),
        ("Docker image: ghcr.io/ise-fizkarlsruhe/ontoink", 0),
        ("", 0),
        ("Ebrahim Norouzi · Tabea Tietz · Harald Sack", 0),
        ("FIZ Karlsruhe — Leibniz Institute for Information Infrastructure", 0),
        ("KIT — Institute AIFB", 0),
    ])

    prs.save("ontoink.pptx")
    return len(prs.slides)


if __name__ == "__main__":
    n = build()
    print(f"Wrote ontoink.pptx ({n} slides)")
